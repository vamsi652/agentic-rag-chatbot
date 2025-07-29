import os
import pandas as pd
import fitz 
from pptx import Presentation
from docx import Document

def parse_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".csv":
        return parse_csv(file_path)
    elif ext in [".txt", ".md"]:
        return parse_txt(file_path)
    else:
        return ""

def parse_pdf(path):
    doc = fitz.open(path)
    text = " ".join(page.get_text() for page in doc)
    doc.close()
    return text

def parse_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return " ".join(text)

def parse_docx(path):
    doc = Document(path)
    return " ".join(p.text for p in doc.paragraphs)

def parse_csv(path):
    df = pd.read_csv(path)
    return df.to_string(index=False)

def parse_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()
